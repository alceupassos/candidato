# -*- coding: utf-8 -*-
"""Gera a apresentação comercial do COCKPIT Eleições 2026 (.pptx)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
PUB = os.path.join(ROOT, "public")

# ── Paleta (do projeto) ──
BG      = RGBColor(0x0F, 0x0F, 0x13)
PANEL   = RGBColor(0x16, 0x16, 0x1D)
CARD    = RGBColor(0x1E, 0x1E, 0x28)
BORDER  = RGBColor(0x2A, 0x2A, 0x38)
VERDE   = RGBColor(0x22, 0xC5, 0x5E)
VERDED  = RGBColor(0x1B, 0x8A, 0x3E)
OURO    = RGBColor(0xF0, 0xC0, 0x30)
AZUL    = RGBColor(0x60, 0xA5, 0xFA)
ROXO    = RGBColor(0x8B, 0x5C, 0xF6)
VERM    = RGBColor(0xEF, 0x44, 0x44)
BRANCO  = RGBColor(0xF0, 0xF0, 0xF0)
SEC     = RGBColor(0x9A, 0x9A, 0xB2)
FONT = "Segoe UI"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H

def img(name, sub="ai"):
    return os.path.join(PUB, sub, name) if sub else os.path.join(PUB, name)

def slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    return s

def rect(s, l, t, w, h, color, line=None, line_w=1.0):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp

def rrect(s, l, t, w, h, color, line=None, line_w=1.0):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    try:
        sp.adjustments[0] = 0.06
    except Exception:
        pass
    return sp

def set_alpha(shape, opacity_pct):
    """Aplica transparência ao preenchimento sólido (opacity_pct: 0–100)."""
    solidFill = shape._element.spPr.find(qn("a:solidFill"))
    if solidFill is None:
        return
    srgb = solidFill.find(qn("a:srgbClr"))
    if srgb is None:
        return
    for old in srgb.findall(qn("a:alpha")):
        srgb.remove(old)
    a = srgb.makeelement(qn("a:alpha"), {"val": str(int(opacity_pct * 1000))})
    srgb.append(a)

def overlay(s, color=RGBColor(0x09, 0x0A, 0x10), opacity=58, l=0, t=0, w=None, h=None):
    sp = rect(s, l, t, w or EMU_W, h or EMU_H, color)
    set_alpha(sp, opacity)
    return sp

def bg(s, color=BG):
    rect(s, 0, 0, EMU_W, EMU_H, color)

def text(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=1.0):
    tb = s.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    first = True
    for line in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align; p.line_spacing = space
        if isinstance(line, tuple):
            line = [line]
        for (txt, size, color, bold) in line:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.name = FONT
    return tb

def pic_fit(s, path, bl, bt, bw, bh, align="center"):
    """Coloca imagem mantendo proporção dentro da caixa."""
    try:
        iw, ih = Image.open(path).size
    except Exception:
        return None
    br = bw / bh; ir = iw / ih
    if ir > br:
        w = bw; h = int(bw / ir)
    else:
        h = bh; w = int(bh * ir)
    left = bl + (bw - w) // 2
    top = bt + (bh - h) // 2
    return s.shapes.add_picture(path, left, top, width=w, height=h)

def pic_cover(s, path, bl, bt, bw, bh):
    """Preenche a caixa (crop) — para fundos."""
    try:
        iw, ih = Image.open(path).size
    except Exception:
        return None
    br = bw / bh; ir = iw / ih
    pic = s.shapes.add_picture(path, bl, bt, width=bw, height=bh)
    if ir > br:
        neww = int(ih * br)
        cl = int((iw - neww) / 2 / iw * 100000)
        pic.crop_left = cl / 100000; pic.crop_right = cl / 100000
    else:
        newh = int(iw / br)
        ct = int((ih - newh) / 2 / ih * 100000)
        pic.crop_top = ct / 100000; pic.crop_bottom = ct / 100000
    return pic

def accent_bar(s, l, t, color=VERDE, w=Inches(0.07), h=Inches(0.5)):
    rect(s, l, t, w, h, color)

def tag(s, l, t, txt, color=VERDE):
    w = Inches(0.18 + 0.092 * len(txt))
    sp = rrect(s, l, t, w, Inches(0.34), PANEL, line=color, line_w=1.0)
    tf = sp.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = txt; r.font.size = Pt(10); r.font.bold = True
    r.font.color.rgb = color; r.font.name = FONT
    return sp

def footer(s, n):
    text(s, Inches(0.5), Inches(7.06), Inches(7), Inches(0.3),
         [[("COCKPIT", 9, VERDE, True), ("  Eleições 2026 · Apresentação confidencial", 9, SEC, False)]])
    text(s, Inches(12.3), Inches(7.06), Inches(0.7), Inches(0.3),
         [[(str(n), 9, SEC, False)]], align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════════════════
# SLIDE 1 — CAPA
# ════════════════════════════════════════════════════════════════════
s = slide(); bg(s)
pic_cover(s, img("hero-warroom.png"), 0, 0, EMU_W, EMU_H)
overlay(s, RGBColor(0x08, 0x0A, 0x0F), 52)            # escurece a foto inteira
overlay(s, RGBColor(0x07, 0x09, 0x0E), 86, w=Inches(8.4))  # aprofunda a esquerda p/ texto

if os.path.exists(img("logovitoriacertabranca.png", sub=None)):
    pic_fit(s, img("logovitoriacertabranca.png", sub=None), Inches(0.7), Inches(0.6), Inches(2.4), Inches(1.0))

accent_bar(s, Inches(0.75), Inches(2.5), VERDE, Inches(0.09), Inches(1.9))
text(s, Inches(1.05), Inches(2.35), Inches(7.2), Inches(2.2), [
    [("COCKPIT", 54, BRANCO, True)],
    [("Eleições ", 54, VERDE, True), ("2026", 54, OURO, True)],
    [("A central de comando da sua campanha.", 20, SEC, False)],
], space=1.05)
tag(s, Inches(1.05), Inches(4.7), "INTELIGÊNCIA ELEITORAL", VERDE)
tag(s, Inches(3.7), Inches(4.7), "TEMPO REAL", OURO)
tag(s, Inches(5.4), Inches(4.7), "IA + DADOS", AZUL)
text(s, Inches(1.05), Inches(5.4), Inches(7), Inches(1.2), [
    [("Transforme energia política em comando, prioridade e ", 15, SEC, False),
     ("voto organizado", 15, VERDE, True), (".", 15, SEC, False)],
    [("Apresentação para diretores de partido e candidatos.", 13, SEC, False)],
], space=1.1)

# ════════════════════════════════════════════════════════════════════
# SLIDE 2 — O CENÁRIO 2026
# ════════════════════════════════════════════════════════════════════
s = slide(); bg(s)
pic_cover(s, img("rally-atmosphere.png"), 0, 0, EMU_W, EMU_H)
overlay(s, RGBColor(0x08, 0x0A, 0x0F), 50)
overlay(s, RGBColor(0x0B, 0x0C, 0x12), 88, w=Inches(8.0))
accent_bar(s, Inches(0.5), Inches(0.6), OURO, Inches(0.09), Inches(0.7))
text(s, Inches(0.8), Inches(0.55), Inches(7), Inches(0.9), [
    [("O desafio de 2026", 30, BRANCO, True)]])
text(s, Inches(0.8), Inches(1.5), Inches(7.0), Inches(5.2), [
    [("Campanhas perdem votos por falta de ", 16, SEC, False), ("comando e dados.", 16, OURO, True)],
    [("", 8, SEC, False)],
    [("• ", 16, VERM, True), ("Decisões no escuro", 16, BRANCO, True), (" — sem saber onde está o voto e o que a população quer.", 15, SEC, False)],
    [("• ", 16, VERM, True), ("Equipe desorganizada", 16, BRANCO, True), (" — líderes, cabos e eleitores sem metas e sem acompanhamento.", 15, SEC, False)],
    [("• ", 16, VERM, True), ("Pesquisas e redes dispersas", 16, BRANCO, True), (" — informação espalhada, lenta e tarde demais.", 15, SEC, False)],
    [("• ", 16, VERM, True), ("Verba e compliance", 16, BRANCO, True), (" — risco com o TSE e gasto sem priorização.", 15, SEC, False)],
    [("• ", 16, VERM, True), ("Mobilização manual", 16, BRANCO, True), (" — WhatsApp e cadastro sem automação nem rastreio.", 15, SEC, False)],
    [("", 10, SEC, False)],
    [("O resultado: esforço alto, conversão baixa.", 17, OURO, True)],
], space=1.18)
footer(s, 2)

# ════════════════════════════════════════════════════════════════════
# SLIDE 3 — A SOLUÇÃO (visão geral)
# ════════════════════════════════════════════════════════════════════
s = slide(); bg(s)
accent_bar(s, Inches(0.5), Inches(0.55), VERDE, Inches(0.09), Inches(0.7))
text(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.9),
     [[("Uma plataforma. Toda a campanha sob controle.", 29, BRANCO, True)]])
text(s, Inches(0.8), Inches(1.35), Inches(11.8), Inches(0.6),
     [[("O COCKPIT reúne inteligência, organização e automação num único painel — no computador e no celular.", 15, SEC, False)]])

cards = [
    ("NOC ao Vivo", "Sala de guerra com pesquisas de Presidente, Senador, Dep. Federal e Estadual rolando em tempo real.", VERDE, "radar"),
    ("Raio-X Regional", "Mostra, por região, o que a população precisa e o que mais converte em voto.", OURO, "scan"),
    ("Meta de Votos", "Funil vivo: meta → inscritos → cadastrados → engajados. Você sabe exatamente onde está.", AZUL, "target"),
    ("Organizadores", "Líderes de Igreja, Gerentes Regionais e Deputados → cabos → eleitores, com metas e monitoramento.", ROXO, "rede"),
    ("Mensageria + IA", "WhatsApp/comunidades, disparos, agente de relacionamento e auto-resposta com IA real.", VERDE, "chat"),
    ("Mapas & Redes", "Mapa do Brasil ▸ estado ▸ região e ranking das redes sociais das principais figuras.", AZUL, "mapa"),
]
cw, ch = Inches(3.95), Inches(2.25)
gx, gy = Inches(0.5), Inches(2.15)
mx, my = Inches(0.18), Inches(0.2)
for i, (t_, d_, c_, _ic) in enumerate(cards):
    col = i % 3; row = i // 3
    l = gx + col * (cw + mx); top = gy + row * (ch + my)
    rrect(s, l, top, cw, ch, CARD, line=BORDER, line_w=1.0)
    rect(s, l, top, Inches(0.08), ch, c_)
    text(s, l + Inches(0.3), top + Inches(0.25), cw - Inches(0.55), Inches(0.5),
         [[(t_, 17, c_, True)]])
    text(s, l + Inches(0.3), top + Inches(0.85), cw - Inches(0.55), Inches(1.3),
         [[(d_, 12.5, RGBColor(0xC9,0xC9,0xDA), False)]], space=1.12)
footer(s, 3)

# ════════════════════════════════════════════════════════════════════
# FEATURE SLIDES
# ════════════════════════════════════════════════════════════════════
def feature(n, titulo, tagtxt, tagcor, subtitulo, bullets, beneficio, image, img_sub="ai", shot=False):
    s = slide(); bg(s)
    # painel de imagem à direita
    panel_l = Inches(7.7)
    rrect(s, panel_l, Inches(0.5), Inches(5.13), Inches(6.0), PANEL, line=BORDER)
    if shot:
        pic_fit(s, image if img_sub is None else img(image, img_sub), panel_l + Inches(0.15), Inches(0.65), Inches(4.83), Inches(5.7))
    else:
        pic_cover(s, image if img_sub is None else img(image, img_sub), panel_l + Inches(0.12), Inches(0.62), Inches(4.89), Inches(5.76))
    # cabeçalho
    accent_bar(s, Inches(0.5), Inches(0.62), tagcor, Inches(0.09), Inches(0.95))
    tag(s, Inches(0.82), Inches(0.55), tagtxt, tagcor)
    text(s, Inches(0.82), Inches(1.0), Inches(6.6), Inches(1.0), [[(titulo, 28, BRANCO, True)]], space=1.0)
    text(s, Inches(0.82), Inches(1.85), Inches(6.6), Inches(0.7), [[(subtitulo, 14, SEC, False)]], space=1.1)
    # bullets
    runs = []
    for b in bullets:
        runs.append([("›  ", 14, tagcor, True), (b[0] + " ", 13.5, BRANCO, True), (b[1], 13, RGBColor(0xC2,0xC2,0xD2), False)])
        runs.append([("", 5, SEC, False)])
    text(s, Inches(0.82), Inches(2.55), Inches(6.6), Inches(3.2), runs, space=1.05)
    # benefício
    by = Inches(5.85)
    rrect(s, Inches(0.82), by, Inches(6.3), Inches(0.95), CARD, line=tagcor, line_w=1.2)
    text(s, Inches(1.05), by + Inches(0.13), Inches(5.9), Inches(0.8),
         [[("BENEFÍCIO  ", 11, tagcor, True), (beneficio, 13.5, BRANCO, True)]], space=1.05, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, n)
    return s

feature(4, "NOC ao Vivo — Sala de Guerra", "MONITORAMENTO", VERDE,
    "War-room com as eleições de Presidente, Senador, Dep. Federal e Estadual em tempo real.",
    [("Abas por cargo", "ranking e evolução de cada disputa, atualizando ao vivo."),
     ("Ticker e alertas", "manchetes e movimentações passando no topo, como uma central de TV."),
     ("Mapa do Brasil", "intenção por estado; clique e mergulhe nas regiões."),
     ("Top oportunidades", "o que mais rende voto, por área e no estado inteiro."),
     ("Motion Report", "vídeo-resumo automático para compartilhar com a equipe.")],
    "Decida com a campanha inteira na tela — sem planilhas, sem atraso.",
    "war-room-campanha.png", img_sub="telas", shot=True)

feature(5, "Raio-X Regional", "INTELIGÊNCIA", OURO,
    "A máquina que mostra o que a população precisa e o que dá mais votos.",
    [("Demanda × Potencial", "cruza o que o eleitor quer com o que converte em voto."),
     ("Índice de oportunidade", "ranking de temas por região, do mais ao menos estratégico."),
     ("Gráficos 3D", "cubo demanda × potencial × satisfação, animado e interativo."),
     ("Mapa de calor", "todas as regiões e temas numa visão só."),
     ("Foco do discurso", "onde falar de saúde, segurança, emprego, turismo…")],
    "Priorize bandeiras que realmente trazem voto — região por região.",
    "ai-eleicoes.png", img_sub="telas", shot=True)

feature(6, "Mapas Drill-Down", "TERRITÓRIO", AZUL,
    "Do Brasil ao bairro: clique e mergulhe na geografia do voto.",
    [("Brasil → Estado → Região", "navegação em camadas, do país ao município."),
     ("Rio de Janeiro como exemplo", "8 regiões de governo e 92 municípios mapeados."),
     ("Cor por desempenho", "cobertura, intenção e oportunidade na cor de cada área."),
     ("Clique filtra tudo", "escolha a região e o painel inteiro se ajusta a ela."),
     ("Vivo e mobile", "atualiza em tempo real e funciona no celular.")],
    "Enxergue onde está o voto e onde investir o próximo comício.",
    "dashboard-politico.png", img_sub="telas", shot=True)

feature(7, "Meta de Votos", "CONVERSÃO", VERDE,
    "O funil que mostra, em tempo real, a distância até a vitória.",
    [("Meta clara", "número-alvo de votos da campanha, por região."),
     ("Funil vivo", "Inscritos → Cadastrados (com dados) → Engajados."),
     ("Engajado = QR", "quem clicou no QR do grupo e participa de pesquisas/mensageria."),
     ("Medidores e gauges", "progresso visual e atualizado a cada interação."),
     ("Dados reais", "soma cadastros, mensagens e respostas de verdade.")],
    "Saiba exatamente quantos votos faltam e de onde virão.",
    "motion-bg.png")

feature(8, "Organizadores de Eleitores", "ESTRUTURA", ROXO,
    "Sua máquina humana organizada em níveis, com metas e monitoramento ao vivo.",
    [("Três lideranças", "Líderes de Igreja, Gerentes Regionais e Deputados Estaduais."),
     ("Base mobilizada", "abaixo deles, cabos eleitorais e eleitores."),
     ("Metas por pessoa", "lista, cadastro (com dados) e engajado (QR)."),
     ("Tempo real", "contatos do dia, conversões, ranking e feed de atividade."),
     ("Controle total", "filtro por nível, organograma e cobrança de meta com 1 clique.")],
    "Cada líder sabe sua meta; você acompanha tudo num só lugar.",
    "community-card.png")

feature(9, "Mensageria + Automação com IA", "RELACIONAMENTO", VERDE,
    "WhatsApp, comunidades e um agente de IA real conversando pela campanha.",
    [("Disparos segmentados", "mensagens por grupo, região e nível de organizador."),
     ("Comunidades e QR", "entrada por QR que conta como engajamento."),
     ("Agente de relacionamento", "responde o eleitor com IA de verdade (não é robô genérico)."),
     ("Templates inteligentes", "personaliza por cidade, bairro e nome automaticamente."),
     ("Pronto para OpenWA", "arquitetura preparada para o WhatsApp oficial.")],
    "Mais conversas, menos esforço — relacionamento que escala.",
    "qr-frame.png")

feature(10, "Ranking de Redes Sociais", "MÍDIA", AZUL,
    "Monitore X, Instagram, YouTube e Facebook das principais figuras políticas.",
    [("Quatro redes", "seguidores, engajamento, sentimento e crescimento."),
     ("Direita, esquerda e centro", "filtre por espectro e por colunistas."),
     ("Buzz por estado", "mapa do Brasil com o volume de menções políticas."),
     ("Benchmark", "compare a presença digital e ache espaços livres."),
     ("Pauta e timing", "saiba quando e sobre o que falar.")],
    "Domine a conversa pública e antecipe o adversário.",
    "noc-banner.png")

feature(11, "Pesquisas, Compliance & Verba", "GESTÃO", OURO,
    "Central de pesquisas, conformidade com o TSE e gestão de verba — tudo integrado.",
    [("Central de pesquisas", "cargo × instituto × região × recorte demográfico."),
     ("Simulador de cenário", "projete resultados ajustando variáveis."),
     ("Compliance TSE", "checklist, prazos e alertas para evitar problemas."),
     ("Gestão de verba", "execução por rubrica, com priorização por oportunidade."),
     ("Calculadora eleitoral", "coeficiente, quociente e metas de votos.")],
    "Decisão técnica, prestação de contas tranquila e gasto inteligente.",
    "login-bg.png")

# ════════════════════════════════════════════════════════════════════
# SLIDE 12 — MOBILE + IA (destaque)
# ════════════════════════════════════════════════════════════════════
s = slide(); bg(s)
pic_cover(s, img("motion-bg.png"), 0, 0, EMU_W, EMU_H)
overlay(s, RGBColor(0x08, 0x0A, 0x0F), 54)
overlay(s, RGBColor(0x0B, 0x0C, 0x12), 84, w=Inches(7.6))
accent_bar(s, Inches(0.5), Inches(0.62), VERDE, Inches(0.09), Inches(0.8))
text(s, Inches(0.82), Inches(0.55), Inches(8), Inches(1.0), [[("War-room no bolso + IA generativa", 28, BRANCO, True)]])
text(s, Inches(0.82), Inches(1.5), Inches(6.8), Inches(5.0), [
    [("Tudo responsivo no celular", 16, VERDE, True)],
    [("Líderes e cabos acompanham metas, recebem tarefas e registram cadastros em campo.", 14, RGBColor(0xC2,0xC2,0xD2), False)],
    [("", 10, SEC, False)],
    [("Inteligência artificial de verdade", 16, OURO, True)],
    [("Conectada via API (Grok, DeepSeek, OpenAI): respostas automáticas, rascunho de mensagens e resumo de notícias.", 14, RGBColor(0xC2,0xC2,0xD2), False)],
    [("", 10, SEC, False)],
    [("Imagens e vídeo sob demanda", 16, AZUL, True)],
    [("Geração de artes e do Motion Report semanal com a identidade da campanha.", 14, RGBColor(0xC2,0xC2,0xD2), False)],
    [("", 10, SEC, False)],
    [("Auto-explicável e animado", 16, ROXO, True)],
    [("Qualquer pessoa da equipe entende em minutos — pensado para vender e para usar.", 14, RGBColor(0xC2,0xC2,0xD2), False)],
], space=1.12)
footer(s, 12)

# ════════════════════════════════════════════════════════════════════
# SLIDE 13 — POR QUE VENCE (diferenciais)
# ════════════════════════════════════════════════════════════════════
s = slide(); bg(s)
accent_bar(s, Inches(0.5), Inches(0.55), OURO, Inches(0.09), Inches(0.7))
text(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.9), [[("Por que o COCKPIT vence", 29, BRANCO, True)]])
difs = [
    ("Tudo em um só lugar", "Inteligência, estrutura e automação integradas — fim das planilhas soltas.", VERDE),
    ("Decisão por dados", "O Raio-X aponta o que dá voto; você para de chutar bandeira.", OURO),
    ("Estrutura sob controle", "Líderes, cabos e eleitores com meta e monitoramento em tempo real.", ROXO),
    ("Automação com IA real", "Relacionamento que escala sem perder o toque humano.", AZUL),
    ("Pronto para o TSE", "Compliance e prestação de contas sem dor de cabeça.", VERDE),
    ("Feito por especialista", "Gestão eleitoral de verdade, não um dashboard genérico.", OURO),
]
cw, ch = Inches(3.95), Inches(2.4)
gx, gy = Inches(0.5), Inches(1.7)
for i, (t_, d_, c_) in enumerate(difs):
    col = i % 3; row = i // 3
    l = gx + col * (cw + Inches(0.18)); top = gy + row * (ch + Inches(0.22))
    rrect(s, l, top, cw, ch, CARD, line=BORDER)
    rect(s, l, top, Inches(0.08), ch, c_)
    text(s, l + Inches(0.3), top + Inches(0.28), cw - Inches(0.55), Inches(0.6), [[(t_, 17, c_, True)]])
    text(s, l + Inches(0.3), top + Inches(0.95), cw - Inches(0.55), Inches(1.3), [[(d_, 13, RGBColor(0xC9,0xC9,0xDA), False)]], space=1.15)
footer(s, 13)

# ════════════════════════════════════════════════════════════════════
# SLIDE 14 — RESULTADOS / IMPACTO
# ════════════════════════════════════════════════════════════════════
s = slide(); bg(s)
pic_cover(s, img("rally-atmosphere.png"), 0, 0, EMU_W, EMU_H)
overlay(s, RGBColor(0x08, 0x0A, 0x0F), 68)
accent_bar(s, Inches(0.5), Inches(0.55), VERDE, Inches(0.09), Inches(0.7))
text(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.9), [[("O que muda na sua campanha", 29, BRANCO, True)]])
metr = [
    ("+ Votos", "organizados e rastreados", VERDE),
    ("- Esforço", "automação no lugar do manual", OURO),
    ("100%", "da estrutura monitorada", AZUL),
    ("Tempo real", "decisão na hora certa", ROXO),
]
cw = Inches(2.9); gx = Inches(0.7); top = Inches(2.5)
for i, (big, sub, c_) in enumerate(metr):
    l = gx + i * (cw + Inches(0.27))
    rrect(s, l, top, cw, Inches(2.4), CARD, line=c_, line_w=1.2)
    text(s, l, top + Inches(0.5), cw, Inches(1.0), [[(big, 36, c_, True)]], align=PP_ALIGN.CENTER)
    text(s, l + Inches(0.2), top + Inches(1.5), cw - Inches(0.4), Inches(0.8), [[(sub, 13, BRANCO, False)]], align=PP_ALIGN.CENTER, space=1.1)
footer(s, 14)

# ════════════════════════════════════════════════════════════════════
# SLIDE 15 — SEGURANÇA & DADOS
# ════════════════════════════════════════════════════════════════════
s = slide(); bg(s)
accent_bar(s, Inches(0.5), Inches(0.55), AZUL, Inches(0.09), Inches(0.7))
text(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.9), [[("Segurança, dados e implantação", 28, BRANCO, True)]])
text(s, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5), [
    [("• ", 16, VERDE, True), ("Acesso protegido", 15, BRANCO, True), (" — login com verificação anti-robô (ALTCHA) e sessão segura.", 14, SEC, False)],
    [("• ", 16, VERDE, True), ("Credenciais por equipe", 15, BRANCO, True), (" — acessos provisórios com validade e limite de dispositivos.", 14, SEC, False)],
    [("• ", 16, VERDE, True), ("Hospedagem dedicada", 15, BRANCO, True), (" — domínio próprio, HTTPS e atualização contínua.", 14, SEC, False)],
    [("• ", 16, VERDE, True), ("Integrações reais e mock", 15, BRANCO, True), (" — IA e imagens ativas; TSE, WhatsApp e redes prontos para conectar.", 14, SEC, False)],
    [("• ", 16, VERDE, True), ("LGPD em mente", 15, BRANCO, True), (" — captação de leads com consentimento e trilha de auditoria.", 14, SEC, False)],
    [("• ", 16, VERDE, True), ("Implantação rápida", 15, BRANCO, True), (" — personalizamos com a identidade e os dados da sua campanha.", 14, SEC, False)],
], space=1.5)
footer(s, 15)

# ════════════════════════════════════════════════════════════════════
# SLIDE 16 — CTA
# ════════════════════════════════════════════════════════════════════
s = slide(); bg(s)
pic_cover(s, img("hero-warroom.png"), 0, 0, EMU_W, EMU_H)
overlay(s, RGBColor(0x07, 0x09, 0x0E), 64)
# overlay via faixa central
ovr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.6), Inches(1.7), Inches(10.1), Inches(4.1))
ovr.fill.solid(); ovr.fill.fore_color.rgb = PANEL; ovr.line.color.rgb = VERDE; ovr.line.width = Pt(1.3); ovr.shadow.inherit=False
text(s, Inches(1.6), Inches(2.05), Inches(10.1), Inches(1.0),
     [[("Coloque a campanha no comando.", 32, BRANCO, True)]], align=PP_ALIGN.CENTER)
text(s, Inches(2.2), Inches(3.0), Inches(8.9), Inches(1.0),
     [[("COCKPIT ", 22, VERDE, True), ("Eleições 2026", 22, OURO, True),
       (" — a vantagem competitiva que separa quem vence de quem assiste.", 18, SEC, False)]],
     align=PP_ALIGN.CENTER, space=1.15)
btn = rrect(s, Inches(4.7), Inches(4.4), Inches(3.9), Inches(0.8), VERDED, line=VERDE, line_w=1.0)
tf = btn.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Agende uma demonstração"; r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = BRANCO; r.font.name = FONT
text(s, Inches(1.6), Inches(5.35), Inches(10.1), Inches(0.4),
     [[("eleicao@angra.io   ·   candidato.angra.io", 13, SEC, False)]], align=PP_ALIGN.CENTER)

OUT = os.path.join(ROOT, "Cockpit-Eleicoes-2026-Apresentacao.pptx")
prs.save(OUT)
print("OK:", OUT, "| slides:", len(prs.slides._sldIdLst))
